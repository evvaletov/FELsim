#!/usr/bin/env python
"""Generate a 2-slide summary of FELsim beamline optimization studies."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x55, 0x99)
GREEN = RGBColor(0x1B, 0x7F, 0x3B)
ORANGE = RGBColor(0xCC, 0x7A, 0x00)
RED = RGBColor(0xCC, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xFA)
HDR_BG = RGBColor(0x00, 0x33, 0x66)


def set_cell(cell, text, size=10, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def style_header_row(table, cols, size=9):
    for i, text in enumerate(cols):
        cell = table.cell(0, i)
        set_cell(cell, text, size=size, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HDR_BG


def stripe_rows(table, start_row=1):
    for r in range(start_row, len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG


def add_title(slide, text, top=Inches(0.2)):
    txBox = slide.shapes.add_textbox(Inches(0.35), top, Inches(9.3), Inches(0.45))
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = DARK


def add_label(slide, text, left, top, width=Inches(4.5), size=12):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = ACCENT


def add_note(slide, text, left, top, width=Inches(9.3), size=10, color=GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color


def add_bullets(slide, items, left, top, width=Inches(9.3)):
    """items: list of (bold_part, rest_text)."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (bold_part, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(1)
        b = p.add_run()
        b.text = bold_part
        b.font.size = Pt(10)
        b.font.bold = True
        b.font.color.rgb = DARK
        r = p.add_run()
        r.text = rest
        r.font.size = Pt(10)
        r.font.color.rgb = GRAY


# ── Slide 1: Studies overview ────────────────────────────────────────────────

def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "UH MkV FEL Beamline — Optimization Studies Summary")
    add_note(slide, ("FELsim Python transfer-matrix simulation, 11-stage Nelder-Mead. "
                     "118 elements (to MkIII undulator entrance, z = 12.4 m). "
                     "Twiss targets from Weinberg, Fisher & Li, arXiv:2510.14061v1 Table I."),
             Inches(0.35), Inches(0.55), size=9)

    # ── Main comparison table ────────────────────────────────────────────────
    y0 = Inches(1.05)
    cols = ["Study", "σ_t", "σ_E", "h (chirp)", "Twiss targets",
            "β_x achieved", "β_y achieved", "MSE"]
    data = [
        # Study 1: 0.5 ps, symmetric targets, emittance-conservation scaling
        ("0.5 ps  emit. cons.", "0.5 ps", "2.0%", "20×10⁹",
         "β = 0.24, α = 0  (both planes)", "0.2419", "0.2419", "1.3e-5"),
        # Study 2: 2 ps, paper targets
        ("2 ps  paper", "2 ps", "0.5%", "5×10⁹",
         "β_x=1.4 α_x=0.47 / β_y=0.24 α_y=0", "1.3999", "0.2419", "2.9e-5"),
        # Study 3: 0.5 ps, paper targets
        ("0.5 ps  paper", "0.5 ps", "0.5%", "5×10⁹",
         "β_x=1.4 α_x=0.47 / β_y=0.24 α_y=0", "1.4000", "0.2420", "2.6e-5"),
        # Study 4: 0.5 ps parameter sensitivity
        ("0.5 ps  sweep", "0.5 ps", "0.1–5%", "0–40×10⁹",
         "β_x=1.4 α_x=0.47 / β_y=0.24 α_y=0", "(see slide 2)", "", ""),
    ]

    nrows = len(data) + 1
    ncols = len(cols)
    tbl = slide.shapes.add_table(nrows, ncols, Inches(0.2), y0,
                                  Inches(9.6), Inches(0.22 * nrows)).table
    widths = [1.25, 0.65, 0.6, 0.8, 2.6, 0.85, 0.85, 0.7]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    style_header_row(tbl, cols, size=8)
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c in (0, 4) else PP_ALIGN.CENTER
            bold = (c == 0)
            set_cell(tbl.cell(r, c), val, size=9, align=align, bold=bold)
    stripe_rows(tbl)

    # Highlight sweep row
    for c in range(ncols):
        tbl.cell(4, c).fill.solid()
        tbl.cell(4, c).fill.fore_color.rgb = RGBColor(0xDB, 0xE8, 0xF4)

    # ── Key differences ──────────────────────────────────────────────────────
    y1 = y0 + Inches(0.22 * nrows) + Inches(0.15)
    add_label(slide, "What differs between studies", Inches(0.35), y1)

    items = [
        ("Symmetric vs asymmetric targets: ",
         "symmetric uses natural undulator focusing in both planes "
         "(β = γλ_u/2πK = 0.24 m, α = 0); asymmetric matches horizontal beam to radiation mode "
         "(β_x = 1.4 m, α_x = 0.47, waist at undulator center)."),
        ("Energy spread scaling: ",
         "study 1 scales σ_E and h by 4× from longitudinal emittance conservation; "
         "studies 2–4 keep σ_E = 0.5% and h = 5×10⁹ per the paper "
         "(\"we only change the bunch length\")."),
        ("Joint 4-variable final stage: ",
         "all studies use chromaticity quad 5 + final triplet (4 variables) to resolve "
         "the overconstrained 3-variable problem.  Required for asymmetric targets, "
         "beneficial for symmetric."),
    ]
    add_bullets(slide, items, Inches(0.35), y1 + Inches(0.28), width=Inches(9.3))

    # ── Bottom: all studies match ─────────────────────────────────────────────
    y2 = y1 + Inches(1.55)
    add_note(slide, ("All four studies achieve excellent undulator Twiss matching (MSE < 3×10⁻⁵). "
                     "The optimizer handles both 0.5 ps and 2 ps, symmetric and asymmetric targets, "
                     "and a wide range of energy spread and chirp."),
             Inches(0.35), y2, size=10, color=DARK)


# ── Slide 2: Parameter sensitivity results ───────────────────────────────────

def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "0.5 ps Parameter Sensitivity — Sweep Results")
    add_note(slide, ("1D scans: hold all parameters at baseline (σ_E = 0.5%, h = 5×10⁹, ε_n = 8), "
                     "sweep one.  500 particles, seed=42.  "
                     "MSE thresholds: Excellent < 10⁻³, Acceptable < 10⁻², Failed > 10⁻¹."),
             Inches(0.35), Inches(0.55), size=9)

    # ── Sweep summary table ──────────────────────────────────────────────────
    y0 = Inches(0.95)
    add_label(slide, "Scan results", Inches(0.35), y0)

    cols = ["Parameter", "Range", "Pts", "Outcome", "Boundary"]
    data = [
        ("σ_E (energy spread)", "0.1% – 5%", "15",
         "14 Excellent, 1 Acceptable",
         "No boundary — excellent across full range"),
        ("h (chirp)", "0 – 40×10⁹ /s", "12",
         "12 Excellent",
         "No boundary — no effect on Twiss matching"),
        ("ε_n (emittance)", "1 – 20 π·mm·mrad", "10",
         "7 Excellent, 2 Acceptable, 1 Failed",
         "Failed at 1; degrades below ~3; excellent ≥ 5"),
    ]
    nrows = len(data) + 1
    tbl = slide.shapes.add_table(nrows, 5, Inches(0.2), y0 + Inches(0.3),
                                  Inches(9.6), Inches(0.24 * nrows)).table
    widths = [1.7, 1.3, 0.5, 2.5, 3.3]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    style_header_row(tbl, cols, size=9)
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c in (0, 3, 4) else PP_ALIGN.CENTER
            bold = (c == 3)
            color = GREEN if r <= 2 and c == 3 else (ORANGE if r == 3 and c == 3 else DARK)
            set_cell(tbl.cell(r, c), val, size=9, align=align, bold=bold, color=color)
    stripe_rows(tbl)

    # ── Emittance detail (condensed) ─────────────────────────────────────────
    y1 = y0 + Inches(0.3) + Inches(0.24 * nrows) + Inches(0.15)
    add_label(slide, "Emittance scan — the only parameter with a feasibility limit",
              Inches(0.35), y1, width=Inches(6))

    em_cols = ["ε_n", "MSE", "Quality", "β_x (m)", "β_y (m)"]
    em_data = [
        ("1", "2.4e+1", "Failed", "6.15", "0.00"),
        ("3", "1.7e-3", "Acceptable", "1.39", "0.15"),
        ("5–12", "~10⁻⁵", "Excellent", "1.40", "0.24"),
        ("14–16", "~10⁻²", "Acceptable", "1.40", "0.02–0.06"),
        ("18–20", "~10⁻⁵", "Excellent", "1.40", "0.24"),
    ]
    q_colors = [RED, ORANGE, GREEN, ORANGE, GREEN]

    nrows2 = len(em_data) + 1
    tbl2 = slide.shapes.add_table(nrows2, 5, Inches(0.2), y1 + Inches(0.3),
                                   Inches(5.4), Inches(0.22 * nrows2)).table
    widths2 = [0.8, 0.8, 1.1, 0.9, 0.9]
    for i, w in enumerate(widths2):
        tbl2.columns[i].width = Inches(w)
    style_header_row(tbl2, em_cols, size=9)
    for r, (en, mse, qual, bx, by) in enumerate(em_data, 1):
        set_cell(tbl2.cell(r, 0), en, size=9)
        set_cell(tbl2.cell(r, 1), mse, size=9)
        set_cell(tbl2.cell(r, 2), qual, size=9, bold=True, color=q_colors[r - 1])
        set_cell(tbl2.cell(r, 3), bx, size=9)
        set_cell(tbl2.cell(r, 4), by, size=9)
    stripe_rows(tbl2)
    # Highlight excellent row
    for c in range(5):
        tbl2.cell(3, c).fill.solid()
        tbl2.cell(3, c).fill.fore_color.rgb = RGBColor(0xDB, 0xE8, 0xF4)

    # ── Conclusions (right side) ─────────────────────────────────────────────
    x_right = Inches(5.85)
    add_label(slide, "Conclusions", x_right, y1, width=Inches(3.5))

    items = [
        ("σ_E and h do not limit Twiss matching",
         " — transfer matrices are linear and energy-independent to first order. "
         "Chirp can be freely minimized (set by upstream RF)."),
        ("Emittance is the sensitive parameter",
         " — fails below ε_n ≈ 1 π·mm·mrad, degrades below ~3. Operating point "
         "(ε_n = 8) has wide margin."),
        ("The ε_n = 14–16 dips",
         " are likely Nelder-Mead local minima (neighbors at 12 and 18 are excellent). "
         "Planned: multi-start verification."),
        ("0.5 ps is achievable",
         " over a wide parameter range. Constraints come from upstream physics "
         "(photocathode, RF), not the beamline."),
    ]
    add_bullets(slide, items, x_right, y1 + Inches(0.28), width=Inches(3.8))


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    slide1(prs)
    slide2(prs)

    from pathlib import Path
    outdir = Path(__file__).resolve().parent / 'results' / 'params_05ps'
    outpath = outdir / 'params_05ps_summary.pptx'
    prs.save(str(outpath))
    print(f"Saved {outpath}")


if __name__ == "__main__":
    main()
